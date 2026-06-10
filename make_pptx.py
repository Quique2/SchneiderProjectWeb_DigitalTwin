"""
Genera la presentación SCADA 4.0 + AI para Equipo 3.
Ejecutar: python make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Paleta ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x06, 0x10, 0x1C)   # fondo principal
BG2      = RGBColor(0x0C, 0x16, 0x22)   # panel
BG3      = RGBColor(0x0A, 0x13, 0x20)   # panel oscuro
GREEN    = RGBColor(0x22, 0xC5, 0x5E)   # verde Schneider
GREEN2   = RGBColor(0x3D, 0xCD, 0x58)   # verde claro
BLUE     = RGBColor(0x38, 0xBD, 0xF8)   # azul info
AMBER    = RGBColor(0xFB, 0xBF, 0x24)   # amarillo warning
RED      = RGBColor(0xEF, 0x44, 0x44)   # rojo crítico
TEAL     = RGBColor(0x2D, 0xD4, 0xBF)   # demo/teal
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT     = RGBColor(0xE6, 0xEE, 0xF7)
MUTED    = RGBColor(0x7D, 0x92, 0xA8)
DIM      = RGBColor(0x54, 0x67, 0x7C)
BORDER   = RGBColor(0x1B, 0x2C, 0x40)
ORANGE   = RGBColor(0xD9, 0x77, 0x40)   # acento cobre

W = Inches(13.33)   # 16:9 ancho
H = Inches(7.5)     # 16:9 alto

# ── Helpers ─────────────────────────────────────────────────────────────────
def bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill=BG2, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, x, y, w, h,
          size=Pt(12), color=TEXT, bold=False, italic=False,
          align=PP_ALIGN.LEFT, wrap=True, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb

def tag(slide, text, x, y, fill=GREEN, text_color=RGBColor(0x06,0x10,0x1C), size=Pt(9)):
    w = Inches(len(text)*0.095 + 0.25)
    h = Inches(0.28)
    r = rect(slide, x, y, w, h, fill=fill)
    txbox(slide, text, x+Inches(0.07), y+Inches(0.02), w, h,
          size=size, color=text_color, bold=True, align=PP_ALIGN.LEFT)
    return w

def accent_bar(slide, color=GREEN, h_=Inches(0.04)):
    rect(slide, 0, 0, W, h_, fill=color)

def section_title(slide, title, sub='', y=Inches(0.25)):
    txbox(slide, title.upper(), Inches(0.55), y, Inches(12), Inches(0.55),
          size=Pt(26), color=WHITE, bold=True)
    if sub:
        txbox(slide, sub, Inches(0.55), y+Inches(0.52), Inches(12), Inches(0.35),
              size=Pt(13), color=MUTED)

def bullet_block(slide, title, items, x, y, w, h, tag_color=GREEN, tag_txt=None):
    rect(slide, x, y, w, h, fill=BG2, line=BORDER)
    # title bar
    rect(slide, x, y, w, Inches(0.32), fill=BG3)
    txbox(slide, title, x+Inches(0.15), y+Inches(0.04), w-Inches(0.2), Inches(0.3),
          size=Pt(11), color=GREEN2, bold=True)
    if tag_txt:
        tag(slide, tag_txt, x+w-Inches(1.4), y+Inches(0.04), fill=tag_color)
    cy = y+Inches(0.4)
    for item in items:
        txbox(slide, f'• {item}', x+Inches(0.18), cy, w-Inches(0.35), Inches(0.32),
              size=Pt(10.5), color=TEXT)
        cy += Inches(0.32)

def kpi(slide, label, value, x, y, w=Inches(2.4), h=Inches(1.1),
        val_color=GREEN2, val_size=Pt(28)):
    rect(slide, x, y, w, h, fill=BG2, line=BORDER)
    txbox(slide, label.upper(), x+Inches(0.15), y+Inches(0.08), w, Inches(0.24),
          size=Pt(8.5), color=MUTED, bold=True, font='Consolas')
    txbox(slide, value, x+Inches(0.15), y+Inches(0.28), w-Inches(0.2), Inches(0.7),
          size=val_size, color=val_color, bold=True, font='Consolas')

def pill(slide, text, x, y, fill=GREEN2, tc=RGBColor(0x06,0x10,0x1C)):
    w = Inches(max(len(text)*0.095+0.3, 0.9))
    r = rect(slide, x, y, w, Inches(0.3), fill=fill)
    txbox(slide, text, x+Inches(0.1), y+Inches(0.04), w, Inches(0.28),
          size=Pt(9), color=tc, bold=True, align=PP_ALIGN.LEFT)
    return w

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]  # completamente en blanco

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 1 — PORTADA                                                  ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s, GREEN, Inches(0.06))

# Panel lateral verde oscuro
rect(s, 0, 0, Inches(4.6), H, fill=RGBColor(0x03, 0x0A, 0x12))

# Título principal
txbox(s, 'Smart CAFI\nInspection System', Inches(0.45), Inches(1.2), Inches(4.0), Inches(1.8),
      size=Pt(30), color=WHITE, bold=True)
txbox(s, 'with AI + SCADA 4.0', Inches(0.45), Inches(2.9), Inches(4.0), Inches(0.6),
      size=Pt(22), color=GREEN2, bold=True)

txbox(s, 'ITESM MTY  ·  RETO SCHNEIDER ELECTRIC 3.0', Inches(0.45), Inches(3.65),
      Inches(4.2), Inches(0.35), size=Pt(9.5), color=MUTED)
txbox(s, 'Automatización del Proceso de Remachado\ne Inspección de Interruptores CAFI',
      Inches(0.45), Inches(4.0), Inches(4.2), Inches(0.65), size=Pt(11), color=TEXT)

# Separador
rect(s, Inches(0.45), Inches(4.8), Inches(3.6), Inches(0.02), fill=BORDER)

# Equipo
txbox(s, 'EQUIPO 3', Inches(0.45), Inches(4.95), Inches(4.0), Inches(0.3),
      size=Pt(9), color=GREEN, bold=True)
members = [
    'Enrique Amir González H.',
    'Diego Becerra Fuentes',
    'Rodrigo Díaz Arrigunaga',
    'Santiago Ordóñez Ramírez',
]
for i, m in enumerate(members):
    txbox(s, m, Inches(0.45), Inches(5.3)+Inches(i*0.27), Inches(4.0), Inches(0.28),
          size=Pt(11), color=TEXT)

txbox(s, 'Campus Monterrey  ·  Junio 2026', Inches(0.45), Inches(6.85),
      Inches(4.0), Inches(0.3), size=Pt(9), color=DIM)

# Panel derecho — descripción y pilares
rect(s, Inches(5.0), Inches(0.5), Inches(7.8), Inches(6.7), fill=BG2, line=BORDER)
txbox(s, 'SCADA 4.0 — 5 PILARES DE IA INDUSTRIAL', Inches(5.3), Inches(0.7),
      Inches(7.2), Inches(0.35), size=Pt(10), color=GREEN, bold=True)

pillars = [
    ('1 & 5', 'The Watchman',    'Detección de Anomalías\n+ Seguridad Industrial',   AMBER, 'Enrique González'),
    ('2',     'The Mechanic',    'Mantenimiento Predictivo',                          BLUE,  'Diego Becerra'),
    ('3',     'The Optimizer',   'Optimización de Procesos',                          GREEN2,'Rodrigo Díaz'),
    ('4',     'The Coordinator', 'Soporte de Decisiones\n+ Guía SOP',                TEAL, 'Santiago Ordóñez'),
]
py = Inches(1.15)
for num, role, desc, col, name in pillars:
    rect(s, Inches(5.15), py, Inches(7.5), Inches(1.25), fill=BG3, line=BORDER)
    rect(s, Inches(5.15), py, Inches(0.55), Inches(1.25), fill=col)
    txbox(s, num, Inches(5.18), py+Inches(0.38), Inches(0.5), Inches(0.5),
          size=Pt(14), color=RGBColor(0x06,0x10,0x1C), bold=True, align=PP_ALIGN.CENTER)
    txbox(s, role.upper(), Inches(5.82), py+Inches(0.07), Inches(4.5), Inches(0.3),
          size=Pt(9.5), color=col, bold=True)
    txbox(s, desc, Inches(5.82), py+Inches(0.38), Inches(4.2), Inches(0.65),
          size=Pt(11), color=TEXT)
    txbox(s, name, Inches(9.9), py+Inches(0.07), Inches(2.7), Inches(0.3),
          size=Pt(9), color=MUTED, align=PP_ALIGN.RIGHT)
    py += Inches(1.3)

txbox(s, 'Deadline: 9 de Junio 2026', Inches(5.3), Inches(6.55),
      Inches(7.2), Inches(0.3), size=Pt(9), color=DIM)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 2 — CONTEXTO DEL SISTEMA                                     ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
section_title(s, 'Contexto del Sistema', '¿Qué construimos y por qué?')

# 3 columnas
cols = [
    ('EL RETO',
     ['Manipulación robótica del CAFI',
      'Inspección de 5 remaches',
      'PASS / NO PASS automático',
      'Eliminar riesgos para el operador',
      'Ciclo < 5 min/pieza'],
     ORANGE),
    ('NUESTRA SOLUCIÓN',
     ['Lexium Cobot + PLC Modicon M262',
      'HMI de 7" para operador',
      'Gemelo digital en tiempo real',
      'Dashboard SCADA 4.0 con IA',
      'Análisis AI (OpenAI + mock local)'],
     GREEN2),
    ('FLUJO DEL CICLO',
     ['1. Operador inicia desde HMI',
      '2. Cobot toma CAFI del canal',
      '3. Coloca en fixture de remachado',
      '4. Ciclo de remachado simulado',
      '5. Inspección: 5 marcadores',
      '6. HMI: PASA / NO PASA'],
     BLUE),
]
for i, (title, items, color) in enumerate(cols):
    x = Inches(0.35 + i * 4.35)
    bullet_block(s, title, items, x, Inches(1.3), Inches(4.1), Inches(5.8),
                 tag_color=color, tag_txt=None)
    rect(s, x, Inches(1.3), Inches(4.1), Inches(0.04), fill=color)

# Arquitectura de comunicaciones
rect(s, Inches(0.35), Inches(6.55), Inches(12.6), Inches(0.65), fill=BG2, line=BORDER)
txbox(s, 'COMUNICACIONES:  WebSocket wss://gateway → Raspberry Pi (Modbus TCP 10.5.5.100:6502) → Lexium Cobot REAL  ·  EtherNet/IP PLC ↔ Cobot  ·  Frontend React / Three.js en Railway',
      Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.55),
      size=Pt(9.5), color=MUTED, font='Consolas')

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 3 — PILAR 1 & 5: ANOMALY DETECTION + INDUSTRIAL SECURITY    ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s, AMBER)
section_title(s, 'Pilares 1 & 5 — The Watchman', 'Enrique Amir González H.  ·  Detección de Anomalías + Seguridad Industrial')
pill(s, 'PILAR 1: ANOMALY DETECTION', Inches(9.2), Inches(0.22), fill=AMBER, tc=RGBColor(0x06,0x10,0x1C))

# PASS / NO PASS lógica — bloque central
rect(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(2.0), fill=BG2, line=BORDER)
rect(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.32), fill=BG3)
txbox(s, 'LÓGICA PASS / NO PASS — ALGORITMO DE DECISIÓN', Inches(0.5), Inches(1.34),
      Inches(10), Inches(0.28), size=Pt(10), color=AMBER, bold=True)

cases = [
    ('✓ PASS', 'CAFI correcto', 'Colocación OK · 5 remaches detectados · Cobot sin falla · Fixture libre', GREEN2),
    ('✗ NO PASS', 'Misalignment CAFI', 'Posición incorrecta en fixture (I_FIXTURE_1_PRESENT = OFF)', RED),
    ('✗ NO PASS', 'Parada inesperada', 'Cobot E-stop / Protective Stop activo · CB_ESTOP / CB_PSTOP', RED),
    ('✗ NO PASS', 'Obstáculo detectado', 'Zona de operación no libre · Alarma de seguridad activa', AMBER),
]
for i, (verdict, cond, detail, col) in enumerate(cases):
    x = Inches(0.5 + i*3.1)
    rect(s, x, Inches(1.68), Inches(3.0), Inches(1.52), fill=BG3, line=BORDER)
    rect(s, x, Inches(1.68), Inches(3.0), Inches(0.28), fill=col)
    txbox(s, verdict, x+Inches(0.1), Inches(1.68), Inches(2.8), Inches(0.28),
          size=Pt(11), color=RGBColor(0x06,0x10,0x1C), bold=True)
    txbox(s, cond, x+Inches(0.1), Inches(2.0), Inches(2.8), Inches(0.28),
          size=Pt(10.5), color=col, bold=True)
    txbox(s, detail, x+Inches(0.1), Inches(2.3), Inches(2.8), Inches(0.82),
          size=Pt(9), color=MUTED)

# Alarmas inteligentes
bullet_block(s, 'GESTIÓN INTELIGENTE DE ALARMAS', [
    'Agrupación por estación: Cobot · Conveyor · Gripper · PLC I/O',
    'Ciclo de vida: ACTIVE → ACKNOWLEDGED → RESOLVED (no se borra por click)',
    'La condición real sigue → NO se puede marcar resuelta',
    'Severidad: CRITICAL (E-stop) · ALARM (paro protectivo) · WARNING (temp > 50°C)',
    'Historial de mantenimiento: comentarios del operador → correo a admin',
], Inches(0.35), Inches(3.45), Inches(6.0), Inches(3.75), tag_color=AMBER)

# Seguridad Industrial (Pilar 5)
bullet_block(s, 'PILAR 5 — SEGURIDAD INDUSTRIAL', [
    'Comunicación segura: WebSocket sobre TLS (ngrok wss://)',
    'API key de OpenAI SOLO en backend Railway (env var), nunca en frontend',
    'Zona de trabajo monitorizada: CB_COLLISION, CB_SOFTLIMIT en tiempo real',
    'E-stop hardware + soft-stop: detección inmediata en SCADA y HMI',
    'Log de auditoría: toda acción del operador registrada con timestamp',
    'NORMAS: ANSI/RIA R15.06 · ISO 10218 · ISO/TS 15066',
], Inches(6.5), Inches(3.45), Inches(6.45), Inches(3.75), tag_color=ORANGE)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 4 — PILAR 2: PREDICTIVE MAINTENANCE                          ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s, BLUE)
section_title(s, 'Pilar 2 — The Mechanic', 'Diego Becerra Fuentes  ·  Mantenimiento Predictivo')
pill(s, 'PILAR 2: PREDICTIVE MAINTENANCE', Inches(8.6), Inches(0.22), fill=BLUE, tc=RGBColor(0x06,0x10,0x1C))

# KPIs de mantenimiento
kpis = [
    ('Joint máx. temp', '50°C→WARN\n60°C→ALARM\n70°C→CRIT', AMBER),
    ('Conveyor ON-time', '60s→WARN\n120s→ALARM', RED),
    ('Presión gripper', '<5.5bar→WARN\n<5.0bar→ALARM', AMBER),
    ('Fuerza EE', '>20N→WARN\n>40N→ALARM', RED),
]
for i, (lbl, val, col) in enumerate(kpis):
    kpi(s, lbl, val, Inches(0.35+i*3.25), Inches(1.3), w=Inches(3.1), h=Inches(1.2), val_color=col, val_size=Pt(14))

# Indicadores del dashboard
bullet_block(s, 'INDICADORES EN DASHBOARD SCADA', [
    'Contador de ciclos de remachado (total acumulado)',
    'Temperatura J1–J6 en tiempo real (WebSocket → gateway → Raspberry Pi)',
    'Historial 120 puntos (~120s) de temperatura de controlador',
    'Conveyor ON-time en vivo + últimas 10 duraciones de encendido (gráfica de barras)',
    'Corriente media del cobot (avg_current_a del controlador)',
    'Speed magnification % — indicador de carga del robot',
], Inches(0.35), Inches(2.65), Inches(6.1), Inches(4.5), tag_color=BLUE)

# Predicción / lógica
bullet_block(s, 'LÓGICA DE PREDICCIÓN Y ALERTAS', [
    'Línea base: temperatura nominal J1–J6 ≈ 42–50°C en operación estándar',
    'Tendencia ascendente sostenida → alarma preventiva antes de superar umbral',
    'Conveyor encendido >60s continuo → riesgo de sobrecalentamiento de motor',
    'Fuerza end-effector >20N → revisar obstrucción o montaje del gripper',
    'Cobot sin movimiento + enabled=false → revisar cadena de seguridad',
    'AI Copilot: diagnóstico local (heurística) + OpenAI sobre snapshot real',
], Inches(6.6), Inches(2.65), Inches(6.35), Inches(4.5), tag_color=BLUE)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 5 — PILAR 3: PROCESS OPTIMIZATION                            ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s, GREEN2)
section_title(s, 'Pilar 3 — The Optimizer', 'Rodrigo Díaz Arrigunaga  ·  Optimización de Procesos')
pill(s, 'PILAR 3: PROCESS OPTIMIZATION', Inches(8.6), Inches(0.22), fill=GREEN2, tc=RGBColor(0x06,0x10,0x1C))

# Parámetros optimizados
rect(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(1.35), fill=BG2, line=BORDER)
txbox(s, 'PARÁMETROS OPTIMIZADOS — EJEMPLO DEMO', Inches(0.5), Inches(1.35),
      Inches(10), Inches(0.28), size=Pt(10), color=GREEN2, bold=True)
params = [
    ('Cobot speed', '85%', 'Cycle time', '≈48s', 'Riveting delay', '30s', 'Inspection timing', '5s'),
]
txbox(s, 'Speed magnification: 85%   ·   Cycle time: ≈48s nominal (43–53s rango real)   ·   Riveting delay: 30s   ·   Inspection timing: 5s',
      Inches(0.5), Inches(1.68), Inches(12.0), Inches(0.85),
      size=Pt(13), color=GREEN2, bold=True, font='Consolas')

# Indicadores
bullet_block(s, 'INDICADORES DE OPTIMIZACIÓN EN DASHBOARD', [
    'Cycle time en vivo: tiempo del ciclo en curso (N/D en REAL sin FSM)',
    'Promedio 10 ciclos: avg_cycle_time_s (modo DEMO: histórico sintético)',
    'Throughput implícito en contador total de ciclos completados',
    'Energy usage: avg_power_w del controlador (WebSocket real)',
    'Cobot speed: speed_magnification_pct en tiempo real',
    'Riveting timing: controlado por la FSM del PLC Modicon M262',
    'AI optimization: recomendaciones de OpenAI sobre snapshot SCADA',
], Inches(0.35), Inches(2.8), Inches(6.1), Inches(4.35), tag_color=GREEN2)

# Recomendaciones IA de ejemplo
rect(s, Inches(6.6), Inches(2.8), Inches(6.35), Inches(4.35), fill=BG2, line=BORDER)
rect(s, Inches(6.6), Inches(2.8), Inches(6.35), Inches(0.32), fill=BG3)
txbox(s, 'EJEMPLO — RECOMENDACIÓN AI COPILOT', Inches(6.75), Inches(2.84),
      Inches(6.0), Inches(0.28), size=Pt(10), color=GREEN2, bold=True)
ai_text = (
    'DIAGNÓSTICO: Operación nominal\n'
    'Speed magnification: 85% — Cobot habilitado\n'
    'Joint máx. temp: 47°C (dentro de umbral)\n\n'
    'RECOMENDACIÓN:\n'
    '→ Reducir speed a 70% si temperatura J6 supera 50°C\n'
    '→ Ciclo promedio 48s → idle time reducible si conveyor\n'
    '   se activa más temprano en secuencia\n'
    '→ Inspection timing: actuar a 4s si cámara ya confirmó\n\n'
    'PUEDE CONTINUAR OPERANDO: SÍ'
)
txbox(s, ai_text, Inches(6.75), Inches(3.18), Inches(6.1), Inches(3.85),
      size=Pt(10), color=TEXT, font='Consolas')

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 6 — PILAR 4: DECISION SUPPORT + SOP                          ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s, TEAL)
section_title(s, 'Pilar 4 — The Coordinator', 'Santiago Ordóñez Ramírez  ·  Soporte de Decisiones + Guía SOP')
pill(s, 'PILAR 4: DECISION SUPPORT + SOP', Inches(8.8), Inches(0.22), fill=TEAL, tc=RGBColor(0x06,0x10,0x1C))

# SOP messages
rect(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(2.4), fill=BG2, line=BORDER)
rect(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.32), fill=BG3)
txbox(s, 'MENSAJES SOP — EJEMPLOS DE OPERACIÓN CORRECTA E INCORRECTA', Inches(0.5), Inches(1.34),
      Inches(12.0), Inches(0.28), size=Pt(10), color=TEAL, bold=True)

# PASS box
rect(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(1.85), fill=RGBColor(0x05,0x1A,0x0C), line=GREEN2)
txbox(s, 'AI STATUS: PASS', Inches(0.65), Inches(1.78), Inches(5.5), Inches(0.32),
      size=Pt(12), color=GREEN2, bold=True, font='Consolas')
txbox(s, 'CAFI correctamente posicionado.\nRemachado exitoso. 5/5 remaches detectados.\nDespacho autorizado. Continúe con el siguiente ciclo.',
      Inches(0.65), Inches(2.1), Inches(5.5), Inches(1.3),
      size=Pt(10.5), color=TEXT, font='Consolas')

# NO PASS box
rect(s, Inches(6.6), Inches(1.7), Inches(6.2), Inches(1.85), fill=RGBColor(0x1A,0x05,0x05), line=RED)
txbox(s, 'AI STATUS: NO PASS', Inches(6.75), Inches(1.78), Inches(5.9), Inches(0.32),
      size=Pt(12), color=RED, bold=True, font='Consolas')
txbox(s, 'Error de posición detectado.\nCAFI fuera de fixture o remache faltante.\nIntervención del operador requerida.\nRepositionar fixture antes de reiniciar el ciclo.',
      Inches(6.75), Inches(2.1), Inches(5.9), Inches(1.3),
      size=Pt(10.5), color=TEXT, font='Consolas')

# Indicadores del coordinador
bullet_block(s, 'INDICADORES EN DASHBOARD — SOPORTE DE DECISIONES', [
    'Etapa actual del proceso (stage): Pick conveyor · Place fixture · Remachado · Inspección · Place bin',
    'Progreso del ciclo: paso N / total N (barra de progreso visual)',
    'Panel de recomendación IA: probable falla · componente sospechoso',
    'SOP instructions: checklist de pasos para el operador',
    'Recovery guidance: acción específica por tipo de alarma',
    'Confirmación de acción del operador: ACK · RESOLVED con timestamp',
], Inches(0.35), Inches(3.85), Inches(6.1), Inches(3.3), tag_color=TEAL)

# Flujo de decisión
bullet_block(s, 'FLUJO INTEGRADO DEL SISTEMA', [
    '1. Cobot ejecuta ciclo → datos al gateway Raspberry Pi',
    '2. AI valida placement (fixture_1_present · camera_pass/fail)',
    '3. PLC registra remachado → SCADA procesa telemetry.io',
    '4. AI Copilot analiza snapshot → PASS / NO PASS',
    '5. Dashboard muestra resultado + instrucción SOP',
    '6. Operador confirma / comenta → log de mantenimiento → correo admin',
    '7. Evento almacenado en historial SCADA',
], Inches(6.6), Inches(3.85), Inches(6.35), Inches(3.3), tag_color=TEAL)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 7 — ARQUITECTURA DE IA                                        ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
section_title(s, 'Arquitectura de IA', 'Motor SCADA 4.0 · AI Copilot · Pipeline de datos en tiempo real')

# Stack de capas
layers = [
    ('CAPA FÍSICA', 'Lexium Cobot  ·  PLC Modicon M262  ·  Gripper neumático  ·  Sistema de inspección', BG3, MUTED),
    ('CAPA GATEWAY', 'Raspberry Pi — Modbus TCP (10.5.5.100:6502)  ·  WebSocket wss://gateway/ws/cobot', BG3, MUTED),
    ('CAPA SCADA ENGINE', 'createScadaEngine() — REAL-ONLY + modo DEMO  ·  Alarmas lifecycle  ·  Historial acotado', BG2, TEXT),
    ('CAPA IA', 'buildAiContext() → POST /api/scada/ai-diagnose (OpenAI gpt-4o)  ·  Fallback: localMockDiagnosis()', RGBColor(0x06,0x14,0x22), BLUE),
    ('CAPA DASHBOARD', 'React + Three.js  ·  SCADA tab: Overview · I/O · Alarmas · Cobot Health · Tendencias', RGBColor(0x04,0x12,0x1A), GREEN2),
]
for i, (name, detail, bg_col, txt_col) in enumerate(layers):
    y = Inches(1.3 + i*1.08)
    rect(s, Inches(0.35), y, Inches(12.6), Inches(1.0), fill=bg_col, line=BORDER)
    rect(s, Inches(0.35), y, Inches(2.2), Inches(1.0), fill=txt_col if i >= 3 else BORDER)
    tc = RGBColor(0x06,0x10,0x1C) if i >= 3 else txt_col
    txbox(s, name, Inches(0.45), y+Inches(0.3), Inches(2.0), Inches(0.4),
          size=Pt(9), color=tc, bold=True)
    txbox(s, detail, Inches(2.7), y+Inches(0.28), Inches(10.0), Inches(0.55),
          size=Pt(10.5), color=TEXT if i < 3 else txt_col, font='Consolas')
    if i < 4:
        txbox(s, '↓', Inches(6.5), y+Inches(1.0), Inches(0.4), Inches(0.15),
              size=Pt(11), color=DIM, align=PP_ALIGN.CENTER)

# AI Security note
rect(s, Inches(0.35), Inches(6.65), Inches(12.6), Inches(0.55), fill=RGBColor(0x0A,0x1A,0x0C), line=GREEN)
txbox(s, '🔒 SEGURIDAD: La API key de OpenAI vive SOLO en Railway backend (env var Scada_Api_Schneider). '
         'El frontend NUNCA la ve — solo hace POST del contexto SCADA serializado. '
         'Sin backend → fallback a heurística local sin llamadas externas.',
      Inches(0.55), Inches(6.68), Inches(12.2), Inches(0.5),
      size=Pt(9.5), color=GREEN2)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 8 — DASHBOARD SCADA: SECCIONES                               ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
section_title(s, 'Dashboard SCADA 4.0', 'Implementación real — 5 secciones + AI Drawer')

sections_data = [
    ('OVERVIEW',
     ['10 métricas esenciales en rejilla 5×2',
      'Estado de celda: NOMINAL / WARNING / ALARM / CRITICAL',
      'Etapa del ciclo + progreso (modo DEMO)',
      'Alarmas activas + peor severidad',
      'Comunicaciones: Cobot WS · PLC I/O · Backend IA'],
     GREEN2),
    ('I/O REAL',
     ['Matriz de señales: 14 inputs + 8 outputs',
      'Cobot: I_COBOT_MOVING · I_COBOT_READY',
      'PLC: conveyor · fixtures · gripper · cámara · stacklight',
      'Estado: REAL · STALE · NOT_CONNECTED · DEMO',
      'Edad del dato en segundos (ageS)'],
     BLUE),
    ('ALARMAS & MANT.',
     ['Lista completa con ciclo de vida ACTIVE→ACK→RESOLVED',
      'No resolvible si condición real sigue presente',
      'Comentario operador → correo automático a admin',
      'Historial de mantenimiento: ACK · RESOLVE · COMMENT',
      'Log de auditoría con timestamp interno'],
     AMBER),
    ('COBOT HEALTH',
     ['KPIs: joint máx. temp · controlador · speed · corriente',
      'Flags: enabled · power · moving · E-stop · collision',
      'Articulaciones J1–J6: ángulo · temp · corriente',
      'TCP position (x,y,z,rx,ry,rz) en mm/°',
      'Force/Torque end-effector en N/Nm'],
     ORANGE),
    ('TENDENCIAS',
     ['Conveyor ON-time en vivo (línea suave SVG)',
      'Últimos 10 encendidos (barra de duración)',
      'Tiempo de ciclo: últimos 10 (barras)',
      'Temperaturas J1–J6 historial 120 puntos',
      'Temperatura controlador: últimos 120s'],
     TEAL),
]
for i, (title, items, color) in enumerate(sections_data):
    col = i % 3
    row = i // 3
    x = Inches(0.35 + col * 4.35)
    y = Inches(1.3 + row * 2.85)
    w = Inches(4.1)
    h = Inches(2.65) if row == 0 else Inches(2.65)
    bullet_block(s, title, items, x, y, w, h, tag_color=color)
    rect(s, x, y, w, Inches(0.04), fill=color)

# AI Drawer
rect(s, Inches(8.85), Inches(4.15), Inches(4.1), Inches(2.65), fill=BG2, line=TEAL)
rect(s, Inches(8.85), Inches(4.15), Inches(4.1), Inches(0.32), fill=BG3)
txbox(s, 'AI DRAWER — MAINTENANCE COPILOT', Inches(9.0), Inches(4.19),
      Inches(4.0), Inches(0.28), size=Pt(10), color=TEAL, bold=True)
ai_items = [
    'Botón "Analizar con IA" → POST snapshot a OpenAI',
    'Resultado: diagnóstico · componente · acción',
    'Checklist de pasos para el operador',
    '"Usar como mensaje a administración"',
    'Fallback local si backend no responde',
]
cy = Inches(4.55)
for item in ai_items:
    txbox(s, f'• {item}', Inches(9.0), cy, Inches(3.85), Inches(0.32),
          size=Pt(10), color=TEXT)
    cy += Inches(0.32)
rect(s, Inches(8.85), Inches(4.15), Inches(0.04), Inches(2.65), fill=TEAL)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 9 — GEMELO DIGITAL 3D + HMI EN VIVO                          ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
section_title(s, 'Gemelo Digital 3D + HMI en Vivo', 'Visualización 3D + Dashboard SCADA + Operación en tiempo real')

bullet_block(s, 'GEMELO DIGITAL — COBOT EN VIVO (3D)', [
    'Modelo URDF del Lexium Cobot V60 con geometría completa renderizado en Three.js/R3F',
    'Posición de articulaciones J1–J6 sincronizadas con telemetría WebSocket en tiempo real',
    'Conveyor, mesa giratoria (NEMA), sistema de remachado y cámara Datalogic en 3D',
    'Layout configurable: Producción · Inspección · Mantenimiento (offsets de componentes)',
    'Cámara 3D interactiva: órbita, zoom, perspectiva con up-vector Z',
    'Panel flotante HMI Operador: arrastrable · redimensionable · overlay sobre la escena',
], Inches(0.35), Inches(1.3), Inches(6.1), Inches(4.0), tag_color=GREEN)

bullet_block(s, 'HMI OPERADOR INTEGRADO', [
    'SVG 480×272 px replica exacta del Harmony ST6 (HMIST6200)',
    'Tags en tiempo real: 7 DI · 5 DO · stacklight · cámara · contadores',
    'Botones activos: START · STOP · RESUME · RESTART · CONFIRMAR · FINALIZAR',
    'Indicadores LED con glow effect según estado real (ON/OFF/STALE)',
    'Detección de datos obsoletos (>3s sin frame) → indicador STALE',
    'Panel overlay flotante: z-index > teach pendant > labels 3D',
], Inches(6.6), Inches(1.3), Inches(6.35), Inches(4.0), tag_color=ORANGE)

# Footer de links
rect(s, Inches(0.35), Inches(5.5), Inches(12.6), Inches(1.7), fill=BG2, line=BORDER)
txbox(s, 'LINKS DEL SISTEMA', Inches(0.55), Inches(5.55),
      Inches(12.0), Inches(0.28), size=Pt(9), color=GREEN, bold=True)
txbox(s, ('Gemelo Digital:      https://schneider-project-web.up.railway.app\n'
          'SCADA (amigo):       https://digitaltwinwebordo-production.up.railway.app/#scada\n'
          'Repo Raspberry Pi:   https://github.com/Quique2/RaspberryPiGIT  (branch master)\n'
          'Gateway Cobot:       wss://unmoral-shrink-cavalry.ngrok-free.dev/ws/cobot'),
      Inches(0.55), Inches(5.88), Inches(12.0), Inches(1.25),
      size=Pt(9.5), color=MUTED, font='Consolas')

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 10 — RESULTADOS Y VALIDACIÓN                                  ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
section_title(s, 'Resultados y Validación', 'Lo que funcionó · Lo que no · Pruebas del sistema')

# Pruebas
rect(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(1.05), fill=BG2, line=BORDER)
rect(s, Inches(0.35), Inches(1.3), Inches(12.6), Inches(0.32), fill=BG3)
txbox(s, 'PRUEBAS DE VALIDACIÓN DEL SISTEMA', Inches(0.5), Inches(1.34),
      Inches(12.0), Inches(0.28), size=Pt(10), color=GREEN, bold=True)
tests = [
    ('Prueba 1', 'CAFI defectuoso\n(< 5 remaches)', 'NO PASA ✗', RED),
    ('Prueba 2', 'CAFI correcto\n(5 remaches)', 'PASA ✓', GREEN2),
    ('Prueba 3', 'E-stop activado\ndurante ciclo', 'NO PASA ✗\n+ ALARMA', RED),
    ('Prueba 4', 'Cobot desconectado\n(demo mode)', 'DEMO activo\nDatos sint.', TEAL),
]
for i, (name, cond, result, col) in enumerate(tests):
    x = Inches(0.5+i*3.15)
    txbox(s, name, x, Inches(1.68), Inches(3.0), Inches(0.25),
          size=Pt(9), color=MUTED, bold=True)
    txbox(s, cond, x, Inches(1.9), Inches(3.0), Inches(0.38),
          size=Pt(9.5), color=TEXT)
    txbox(s, result, x, Inches(2.25), Inches(3.0), Inches(0.38),
          size=Pt(10), color=col, bold=True, font='Consolas')

# Lo que funcionó
bullet_block(s, 'LO QUE FUNCIONÓ BIEN ✓', [
    'Gemelo digital 3D sincronizado con cobot REAL vía WebSocket + Raspberry Pi',
    'HMI Operador SVG fiel al Harmony ST6 — señales en tiempo real',
    'Motor SCADA REAL-ONLY: alarmas con ciclo de vida completo',
    'AI Copilot: OpenAI + fallback local — nunca falla por falta de backend',
    'Deploy continuo en Railway — workflow git commit → producción < 2 min',
    'Modo DEMO: permite demostrar el sistema sin hardware conectado',
], Inches(0.35), Inches(2.75), Inches(6.1), Inches(4.4), tag_color=GREEN2)

# Lo que no funcionó
bullet_block(s, 'ÁREAS DE MEJORA / DIFICULTADES ✗', [
    'PLC I/O (telemetry.io): el gateway actual NO publica el bloque io del PLC',
    '→ Conveyor, fixtures, gripper, cámara aparecen como NOT_CONNECTED',
    'Ciclo de tiempo REAL: SCADA aislado de la FSM → cycleStep = N/D en REAL',
    'HMI 3D embebida en pantalla del modelo: transform-fixed no visible a distancia',
    'Tiempo de ciclo real no integrado: no hay bloque de ciclo del PLC en WebSocket',
    'Latencia ngrok: reconexión WebSocket cada ~4s si el túnel se interrumpe',
], Inches(6.6), Inches(2.75), Inches(6.35), Inches(4.4), tag_color=AMBER)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 11 — MEJORAS Y CONCLUSIONES                                   ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s, GREEN)
section_title(s, 'Mejoras y Conclusiones', 'Trabajo futuro · Reflexión del equipo · Lecciones aprendidas')

bullet_block(s, 'MEJORAS PROPUESTAS — SIGUIENTE ITERACIÓN', [
    'Publicar bloque telemetry.io en el gateway RPi → activar ALL las señales PLC en SCADA',
    'Integrar bloque de ciclo (stage, step, total) en el frame WebSocket → cycleStep REAL',
    'Agregar MQTT broker para PLC Modicon M262 → mayor robustez que polling Modbus',
    'Implementar TensorFlow Lite en Raspberry Pi para anomaly detection EDGE',
    'HMI 3D: escalar correctamente con la pantalla física del modelo en la celda',
    'ML predictivo: regresión sobre historial de temperaturas para RUL de articulaciones',
    'Persistent storage: PostgreSQL en Railway para historial largo de alarmas y ciclos',
], Inches(0.35), Inches(1.3), Inches(6.1), Inches(5.0), tag_color=GREEN)

bullet_block(s, 'CONCLUSIONES Y REFLEXIÓN DEL EQUIPO', [
    'Logramos un gemelo digital funcional conectado al cobot REAL via Modbus TCP',
    'El SCADA 4.0 cubre los 5 pilares IA con datos reales y modo demo',
    'La integración OpenAI permite diagnóstico inteligente sin exponer credenciales',
    'El HMI SVG replica fielmente el panel físico Harmony ST6 del laboratorio',
    'El sistema es extensible: cada pilar puede incorporar ML sin afectar los demás',
    'Detectamos: separar SCADA de FSM fue correcto — evita acoplamiento frágil',
    'Lo mejor: arquitectura real → física RPi → cobot → WebSocket → React en < 200ms',
], Inches(6.6), Inches(1.3), Inches(6.35), Inches(5.0), tag_color=TEAL)

# Criterios de éxito finales
rect(s, Inches(0.35), Inches(6.5), Inches(12.6), Inches(0.7), fill=RGBColor(0x03,0x12,0x08), line=GREEN)
checks = [
    ('✓ Cobot manipula el CAFI', GREEN2),
    ('✓ PASA / NO PASA automático', GREEN2),
    ('✓ HMI Operador REAL', GREEN2),
    ('✓ SCADA 4.0 con 5 pilares IA', GREEN2),
    ('⚠ telemetry.io PLC pendiente', AMBER),
]
for i, (txt, col) in enumerate(checks):
    txbox(s, txt, Inches(0.55+i*2.5), Inches(6.58), Inches(2.4), Inches(0.5),
          size=Pt(10), color=col, bold=True)

# ╔══════════════════════════════════════════════════════════════════════╗
# ║  SLIDE 12 — RÚBRICA DE EVALUACIÓN                                    ║
# ╚══════════════════════════════════════════════════════════════════════╝
s = prs.slides.add_slide(blank_layout)
bg(s)
accent_bar(s)
section_title(s, 'Autoevaluación — Rúbrica', 'Criterios de evaluación cubiertos por nuestra implementación')

criteria = [
    ('1. Data Analysis &\nInterpretation', '25 pts',
     'Motor SCADA REAL-ONLY: deriva estado SOLO de telemetría real. '
     'Identifica patrones: temperatura ascendente, conveyor on-time, '
     'fuerza EE, flags cobot. Historial acotado (120 puntos) para tendencias.',
     GREEN2, Inches(0.35)),
    ('2. AI-Driven Insights &\nRecommendations', '20 pts',
     'AI Copilot: buildAiContext() → OpenAI gpt-4o con snapshot SCADA completo. '
     'Diagnóstico probable · componente sospechoso · acción · checklist. '
     'Fallback heurístico local. Modo DEMO marcado explícitamente para la IA.',
     BLUE, Inches(0.35)),
    ('3. Data Visualization\nQuality', '20 pts',
     'Dashboard industrial brutalista: 5 secciones, gráficas SVG (barras + líneas '
     'Catmull-Rom). Gemelo 3D con Three.js/R3F. HMI SVG 480×272. '
     'Color SOLO para estado (semáforo industrial). Layout max-width 1440.',
     AMBER, Inches(0.35)),
    ('4. Presentation &\nCommunication', '15 pts',
     'Presentación estructurada por pilares SCADA 4.0. Evidencia: '
     'URLs de Railway en producción. Nomenclatura industrial correcta. '
     'PASS/NO PASS con lógica explícita y ejemplos.',
     ORANGE, Inches(0.35)),
    ('5. Teamwork &\nProject Deliverable', '20 pts',
     'Enrique (Watchman 1&5) · Diego (Mechanic 2) · Rodrigo (Optimizer 3) '
     '· Santiago (Coordinator 4). Sistema funcional entregado antes del 9 Jun 2026. '
     'GitHub + Railway + RPi integrados.',
     TEAL, Inches(0.35)),
]
for i, (crit, pts, desc, col, _) in enumerate(criteria):
    row = i // 2 if i < 4 else None
    col_idx = i % 2
    if i < 4:
        x = Inches(0.35 + col_idx * 6.5)
        y = Inches(1.3 + (i // 2) * 2.6)
        w = Inches(6.2)
        h = Inches(2.4)
    else:
        x = Inches(0.35)
        y = Inches(6.5)
        w = Inches(12.6)
        h = Inches(0.75)

    rect(s, x, y, w, h, fill=BG2, line=BORDER)
    rect(s, x, y, Inches(0.06), h, fill=col)
    txbox(s, crit, x+Inches(0.15), y+Inches(0.08), w-Inches(0.6), Inches(0.5),
          size=Pt(11), color=col, bold=True)
    txbox(s, pts, x+w-Inches(0.9), y+Inches(0.08), Inches(0.8), Inches(0.35),
          size=Pt(14), color=col, bold=True, align=PP_ALIGN.RIGHT, font='Consolas')
    txbox(s, desc, x+Inches(0.15), y+Inches(0.58 if i < 4 else 0.18),
          w-Inches(0.25), h-Inches(0.65 if i < 4 else 0.25),
          size=Pt(9.5), color=TEXT)

out = r'c:\Users\kiki7\Downloads\Schneider Project\SCADA_AI_Equipo3.pptx'
prs.save(out)
print(f'OK Guardado: {out}')
